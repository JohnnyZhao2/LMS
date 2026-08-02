"""文档解析：DOCX/PPTX/PDF → HTML，及唯一解析接口。"""

from __future__ import annotations

import os
import re

from django.utils.html import escape
from drf_spectacular.utils import extend_schema
from rest_framework.parsers import MultiPartParser
from rest_framework.permissions import IsAuthenticated
from rest_framework.views import APIView

from apps.authorization.engine import enforce
from core.exceptions import BusinessError, ErrorCodes
from core.responses import success_response


class DocumentParserService:
    """文档解析服务。"""

    SUPPORTED_EXTENSIONS = {'.docx', '.pptx', '.pdf'}
    MAX_FILE_SIZE = 10 * 1024 * 1024
    DOCX_HEADING_STYLE_PATTERNS = (
        re.compile(r'^Heading\s*(\d+)$', re.IGNORECASE),
        re.compile(r'^标题\s*(\d+)$'),
    )
    DOCX_DECIMAL_HEADING_PATTERN = re.compile(r'^\s*(\d+(?:[.．]\d+)+)\s*\S+')
    INLINE_ORDERED_LIST_PATTERN = re.compile(r'(?<![\d.．])(\d+)[.．、]\s+')
    DOCX_TITLE_STYLES = {'Title', '标题'}
    DOCX_SUBTITLE_STYLES = {'Subtitle', '副标题'}

    def parse(self, file) -> tuple[str, str]:
        if file.size > self.MAX_FILE_SIZE:
            raise BusinessError(
                code=ErrorCodes.VALIDATION_ERROR,
                message=f'文件大小超过限制（最大 {self.MAX_FILE_SIZE // 1024 // 1024}MB）',
            )

        filename = file.name
        ext = os.path.splitext(filename)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise BusinessError(
                code=ErrorCodes.VALIDATION_ERROR,
                message=f'不支持的文件格式，仅支持 {", ".join(sorted(self.SUPPORTED_EXTENSIONS))}',
            )

        if ext == '.docx':
            return self._parse_docx(file)
        if ext == '.pptx':
            return self._parse_pptx(file)
        return self._parse_pdf(file)

    def _parse_docx(self, file) -> tuple[str, str]:
        from docx import Document

        doc = Document(file)
        html_parts = []
        title = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                html_parts.append('<p><br></p>')
                continue
            style_name = para.style.name if para.style else ''
            heading_level = self._resolve_docx_heading_level(para, text)
            if heading_level:
                if title is None and heading_level == 1:
                    title = text
                content_level = self._resolve_docx_content_heading_level(heading_level)
                html_parts.append(f'<h{content_level}>{escape(text)}</h{content_level}>')
            elif style_name == 'List Bullet':
                html_parts.append(f'<ul><li>{escape(text)}</li></ul>')
            elif style_name == 'List Number':
                html_parts.append(f'<ol><li>{escape(text)}</li></ol>')
            else:
                html_parts.extend(self._render_text_blocks(text))

        content = '\n'.join(html_parts)
        content = self._merge_consecutive_lists(content)
        return title or self._extract_title_from_filename(file.name), content

    def _resolve_docx_heading_level(self, para, text: str) -> int | None:
        style_name = para.style.name if para.style else ''

        numbered_level = self._resolve_docx_numbered_heading_level(text)
        if numbered_level:
            return numbered_level

        style_level = self._resolve_docx_style_heading_level(style_name)
        if style_level:
            return style_level

        outline_level = self._resolve_docx_outline_level(para)
        if outline_level:
            return outline_level

        return None

    def _resolve_docx_numbered_heading_level(self, text: str) -> int | None:
        level_match = self.DOCX_DECIMAL_HEADING_PATTERN.match(text.strip())
        if not level_match:
            return None

        section_number = level_match.group(1).replace('．', '.')
        return min(section_number.count('.') + 1, 6)

    def _resolve_docx_content_heading_level(self, heading_level: int) -> int:
        return min(heading_level + 1, 6)

    def _resolve_docx_style_heading_level(self, style_name: str) -> int | None:
        normalized_style_name = style_name.strip()
        if normalized_style_name in self.DOCX_TITLE_STYLES:
            return 1
        if normalized_style_name in self.DOCX_SUBTITLE_STYLES:
            return 2

        for pattern in self.DOCX_HEADING_STYLE_PATTERNS:
            level_match = pattern.match(normalized_style_name)
            if level_match:
                return int(level_match.group(1))
        return None

    def _resolve_docx_outline_level(self, para) -> int | None:
        direct_level = self._read_docx_outline_level_value(getattr(getattr(para, '_p', None), 'pPr', None))
        if direct_level:
            return direct_level

        style_element = getattr(getattr(para, 'style', None), 'element', None)
        return self._read_docx_outline_level_value(getattr(style_element, 'pPr', None))

    def _read_docx_outline_level_value(self, ppr) -> int | None:
        outline_level = getattr(ppr, 'outlineLvl', None)
        if outline_level is None:
            return None

        raw_value = outline_level.val
        if isinstance(raw_value, int):
            return raw_value + 1
        if isinstance(raw_value, str) and raw_value.isdigit():
            return int(raw_value) + 1
        return None

    def _parse_pptx(self, file) -> tuple[str, str]:
        from pptx import Presentation

        prs = Presentation(file)
        html_parts = []
        title = None

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            slide_title = None
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    if slide_title is None:
                        slide_title = text
                    else:
                        slide_texts.append(text)
            if slide_title:
                if title is None and i == 1:
                    title = slide_title
                html_parts.append(f'<h2>第 {i} 页：{escape(slide_title)}</h2>')
                for text in slide_texts:
                    for line in text.split('\n'):
                        line = line.strip()
                        if line:
                            html_parts.extend(self._render_text_blocks(line))
                        else:
                            html_parts.append('<p><br></p>')

        content = '\n'.join(html_parts)
        content = self._merge_consecutive_lists(content)
        return title or self._extract_title_from_filename(file.name), content

    def _parse_pdf(self, file) -> tuple[str, str]:
        import pdfplumber

        html_parts = []
        title = None
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    text = text.strip()
                    lines = text.split('\n')
                    if title is None and i == 1 and lines:
                        title = lines[0].strip()
                    html_parts.append(f'<h2>第 {i} 页</h2>')
                    for line in lines:
                        line = line.strip()
                        if line:
                            html_parts.extend(self._render_text_blocks(line))
                        else:
                            html_parts.append('<p><br></p>')
        content = '\n'.join(html_parts)
        content = self._merge_consecutive_lists(content)
        return title or self._extract_title_from_filename(file.name), content

    def _render_text_blocks(self, text: str) -> list[str]:
        inline_list = self._split_inline_ordered_list(text)
        if not inline_list:
            return [f'<p>{escape(text)}</p>']

        prefix, items = inline_list
        blocks = []
        if prefix:
            blocks.append(f'<p>{escape(prefix)}</p>')
        blocks.append('<ol>' + ''.join(f'<li>{escape(item)}</li>' for item in items) + '</ol>')
        return blocks

    def _split_inline_ordered_list(self, text: str) -> tuple[str, list[str]] | None:
        matches = list(self.INLINE_ORDERED_LIST_PATTERN.finditer(text))
        if not matches:
            return None

        first_match = matches[0]
        prefix = text[:first_match.start()].strip()
        if prefix and (first_match.group(1) != '1' or len(matches) == 1):
            return None

        items = []
        for index, match in enumerate(matches):
            next_match = matches[index + 1] if index + 1 < len(matches) else None
            item = text[match.end():next_match.start() if next_match else len(text)].strip()
            item = re.sub(r'[;；]\s*$', '', item).strip()
            if item:
                items.append(item)

        return (prefix, items) if items else None

    def _merge_consecutive_lists(self, html: str) -> str:
        html = re.sub(r'</ul>\s*<ul>', '', html)
        html = re.sub(r'</ol>\s*<ol>', '', html)
        return html

    def _extract_title_from_filename(self, filename: str) -> str:
        return os.path.splitext(filename)[0] or '未命名文档'


class ParseDocumentView(APIView):
    """POST /api/knowledge/parse-document/ — 上传文档并解析为 HTML。"""

    parser_classes = [MultiPartParser]
    permission_classes = [IsAuthenticated]

    @extend_schema(
        summary='解析上传文档',
        description='上传 DOCX/PPTX/PDF，解析为 HTML，供创建知识前预填。',
        tags=['知识管理'],
    )
    def post(self, request):
        # capability gate：仅创建知识权限可解析
        enforce('knowledge.create', request, error_message='无权解析文档')
        file = request.FILES.get('file')
        if not file:
            raise BusinessError(code=ErrorCodes.VALIDATION_ERROR, message='请上传文件')

        suggested_title, content = DocumentParserService().parse(file)
        ext = file.name.rsplit('.', 1)[-1].lower() if '.' in file.name else ''
        return success_response({
            'suggested_title': suggested_title,
            'content': content,
            'file_type': ext,
        })

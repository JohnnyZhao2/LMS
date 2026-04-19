"""Knowledge services."""

import hashlib
import json
import os
import re
from typing import Optional, Tuple

from apps.activity_logs.decorators import log_content_action
from django.db import transaction
from django.utils.html import escape, strip_tags

from apps.tags.resource_sync import (
    apply_resource_tag_changes,
    build_resource_update_plan,
    pop_resource_tag_payload,
)
from core.base_service import BaseService
from core.exceptions import BusinessError, ErrorCodes

from .models import Knowledge, KnowledgeRevision
from .selectors import get_knowledge_by_id


def build_knowledge_revision_payload(knowledge: Knowledge) -> dict:
    return {
        'title': knowledge.title,
        'content': knowledge.content,
        'related_links': knowledge.related_links,
        'space_tag_name': knowledge.space_tag.name if knowledge.space_tag else '',
        'tags_json': [
            {'id': tag.id, 'name': tag.name, 'tag_type': tag.tag_type}
            for tag in knowledge.tags.order_by('id')
        ],
    }


def build_knowledge_revision_hash(payload: dict) -> str:
    normalized = json.dumps(payload, sort_keys=True, ensure_ascii=False, default=str)
    return hashlib.sha256(normalized.encode('utf-8')).hexdigest()


def ensure_knowledge_revision(knowledge: Knowledge, *, actor) -> KnowledgeRevision:
    payload = build_knowledge_revision_payload(knowledge)
    content_hash = build_knowledge_revision_hash(payload)
    latest = KnowledgeRevision.objects.filter(source_knowledge=knowledge).order_by('-revision_number').first()
    if latest and latest.content_hash == content_hash:
        return latest

    next_revision_number = (latest.revision_number if latest else 0) + 1
    return KnowledgeRevision.objects.create(
        source_knowledge=knowledge,
        revision_number=next_revision_number,
        title=payload['title'],
        content=payload['content'],
        related_links=payload['related_links'],
        space_tag_name=payload['space_tag_name'],
        tags_json=payload['tags_json'],
        content_hash=content_hash,
        created_by=actor,
    )


class KnowledgeService(BaseService):
    """知识文档应用服务。"""

    def get_by_id(self, pk: int) -> Knowledge:
        knowledge = get_knowledge_by_id(pk)
        self.validate_not_none(knowledge, f'知识文档 {pk} 不存在')
        return knowledge

    @transaction.atomic
    @log_content_action(
        'knowledge',
        'create',
        '',
        group='知识文档',
        label='创建知识文档',
    )
    def create(self, data: dict) -> Knowledge:
        self._validate_knowledge_data(data)
        payload = dict(data)
        payload['created_by'] = self.user
        payload['updated_by'] = self.user
        tag_payload = pop_resource_tag_payload(payload, scope='knowledge')
        knowledge = Knowledge.objects.create(**payload)
        apply_resource_tag_changes(
            knowledge,
            space_tag_id=tag_payload.space_tag_id,
            tag_ids=tag_payload.tag_ids,
            space_tag_provided=tag_payload.space_tag_provided,
            tag_ids_provided=True,
        )
        return knowledge

    @transaction.atomic
    @log_content_action(
        'knowledge',
        'update',
        '',
        group='知识文档',
        label='更新知识文档',
    )
    def update(self, pk: int, data: dict) -> Knowledge:
        knowledge = self.get_by_id(pk)
        self._validate_knowledge_data(data=data, fallback_content=knowledge.content)

        current_tag_ids = list(knowledge.tags.values_list('id', flat=True))
        update_plan = build_resource_update_plan(
            knowledge,
            data,
            scope='knowledge',
            current_tag_ids=current_tag_ids,
        )
        if not update_plan.has_changes:
            return knowledge

        changed_fields = dict(update_plan.changed_fields)
        changed_fields['updated_by'] = self.user
        for key, value in changed_fields.items():
            setattr(knowledge, key, value)
        knowledge.save(update_fields=list(changed_fields.keys()))
        apply_resource_tag_changes(
            knowledge,
            space_tag_id=update_plan.space_tag_id,
            tag_ids=update_plan.tag_ids,
            space_tag_provided=update_plan.space_changed,
            tag_ids_provided=update_plan.tags_changed,
        )
        return knowledge

    @transaction.atomic
    @log_content_action(
        'knowledge',
        'delete',
        '',
        group='知识文档',
        label='删除知识文档',
    )
    def delete(self, pk: int) -> Knowledge:
        knowledge = self.get_by_id(pk)
        revisions = list(knowledge.revisions.all())
        knowledge.delete()
        KnowledgeRevision.objects.filter(
            id__in=[revision.id for revision in revisions],
            knowledge_tasks__isnull=True,
        ).delete()
        return knowledge

    def increment_view_count(self, pk: int) -> int:
        return self.get_by_id(pk).increment_view_count()

    def _validate_knowledge_data(
        self,
        data: dict,
        fallback_content: Optional[str] = None,
    ) -> None:
        effective_content = data.get('content', fallback_content or '')
        if not strip_tags(str(effective_content)).strip():
            raise BusinessError(
                code=ErrorCodes.VALIDATION_ERROR,
                message='知识文档必须填写正文内容',
            )


class DocumentParserService:
    """文档解析服务。"""

    SUPPORTED_EXTENSIONS = {'.docx', '.pptx', '.pdf'}
    MAX_FILE_SIZE = 10 * 1024 * 1024

    def parse(self, file) -> Tuple[str, str]:
        if file.size > self.MAX_FILE_SIZE:
            raise ValueError(f'文件大小超过限制（最大 {self.MAX_FILE_SIZE // 1024 // 1024}MB）')

        filename = file.name
        ext = os.path.splitext(filename)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f'不支持的文件格式，仅支持 {", ".join(self.SUPPORTED_EXTENSIONS)}')

        if ext == '.docx':
            return self._parse_docx(file)
        if ext == '.pptx':
            return self._parse_pptx(file)
        return self._parse_pdf(file)

    def _parse_docx(self, file) -> Tuple[str, str]:
        from docx import Document

        doc = Document(file)
        html_parts = []
        title = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style else ''
            if style_name.startswith('Heading'):
                level_match = re.search(r'\d+', style_name)
                level = int(level_match.group()) if level_match else 1
                level = min(level, 6)
                if title is None and level == 1:
                    title = text
                html_parts.append(f'<h{level}>{escape(text)}</h{level}>')
            elif style_name == 'List Bullet':
                html_parts.append(f'<ul><li>{escape(text)}</li></ul>')
            elif style_name == 'List Number':
                html_parts.append(f'<ol><li>{escape(text)}</li></ol>')
            else:
                html_parts.append(f'<p>{escape(text)}</p>')

        content = '\n'.join(html_parts)
        content = self._merge_consecutive_lists(content)
        return title or self._extract_title_from_filename(file.name), content

    def _parse_pptx(self, file) -> Tuple[str, str]:
        from pptx import Presentation

        prs = Presentation(file)
        html_parts = []
        title = None

        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            slide_title = None
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    if slide_title is None:
                        slide_title = text
                    else:
                        slide_texts.append(text)
            if slide_title:
                if title is None and i == 1:
                    title = slide_title
                html_parts.append(f'<h2>第 {i} 页：{escape(slide_title)}</h2>')
                for text in slide_texts:
                    for line in text.split('\n'):
                        line = line.strip()
                        if line:
                            html_parts.append(f'<p>{escape(line)}</p>')

        return title or self._extract_title_from_filename(file.name), '\n'.join(html_parts)

    def _parse_pdf(self, file) -> Tuple[str, str]:
        import pdfplumber

        html_parts = []
        title = None
        with pdfplumber.open(file) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    text = text.strip()
                    lines = text.split('\n')
                    if title is None and i == 1 and lines:
                        title = lines[0].strip()
                    html_parts.append(f'<h2>第 {i} 页</h2>')
                    for line in lines:
                        line = line.strip()
                        if line:
                            html_parts.append(f'<p>{escape(line)}</p>')
        return title or self._extract_title_from_filename(file.name), '\n'.join(html_parts)

    def _merge_consecutive_lists(self, html: str) -> str:
        html = re.sub(r'</ul>\s*<ul>', '', html)
        html = re.sub(r'</ol>\s*<ol>', '', html)
        return html

    def _extract_title_from_filename(self, filename: str) -> str:
        return os.path.splitext(filename)[0] or '未命名文档'
